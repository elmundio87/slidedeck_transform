import argparse
import sys
import json
import os
import copy
from pptx import Presentation
from PIL import Image, ImageChops
import regex


TEMP_LOGO_FILE = "temp_logo.png"


def uniq(items):
    return list(set(items))


def is_json(text):
    try:
        json.loads(text)
    except ValueError:
        return False
    return True


def get_all_tags_in_presentation(presentation):
    tags = []
    for slide in presentation.slides:
        if slide.has_notes_slide:
            text_frame = slide.notes_slide.notes_text_frame
            tags += get_all_tags_in_comment(text_frame.text)
    return uniq(tags)


def find_json_in_string(text):
    return regex.findall('{(?:[^{}]|(?R))*}', text)


def get_all_tags_in_comment(text):
    json_list = find_json_in_string(text)
    all_tags = []
    if json_list:
        for json_block in json_list:
            if is_json(json_block):
                tags = json.loads(json_block)['tags']
                if not isinstance(tags, list):
                    tags = [tags]
                all_tags += tags
    return uniq(all_tags)


def remove_json_from_comment(text):
    json_list = find_json_in_string(text)
    if json_list:
        for json_block in json_list:
            text = text.replace(json_block, "")
    return text


def trim(image):
    background = Image.new(image.mode, image.size, image.getpixel((0, 0)))
    diff = ImageChops.difference(image, background)
    diff = ImageChops.add(diff, diff, 1.0, -100)
    bbox = diff.getbbox()
    if bbox:
        return image.crop(bbox)


def delete_slide(prs, slide):
    # Make dictionary with necessary information
    id_dict = {slide.id: [i, slide.rId]
               for i, slide in enumerate(prs.slides._sldIdLst)}
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]


def parse_args():
    parser = argparse.ArgumentParser(
        description='Transform a pptx presentation')
    required = parser.add_argument_group('Required arguments')
    optional = parser.add_argument_group('Optional arguments')
    optional.add_argument('--trim-logo', action='store_true',
                          help="Trim any whitespace from the logo")
    required.add_argument('--file', nargs=1, action='store',
                          help="Path to the powerpoint presentation file")
    required.add_argument('--logo', nargs=1, action='store',
                          help="Path to the logo file")
    optional.add_argument('--tags', nargs="?", action='store', default="",
                          help="List of tags to remove")

    return parser.parse_args()


def main():
    ARGS = parse_args()

    print 'Argument List:', ARGS
    presentation_file = ARGS.file[0]

    presentation = Presentation(presentation_file)

    logo_image = ARGS.logo[0]
    # logo_image = "/Users/edmundd/Desktop/admiral_logo.png"
    # logo_image = "/Users/edmundd/Desktop/logo_gb.png"
    # logo_image = "/Users/edmundd/Desktop/Nokia-logo.jpg"
    # logo_image = "/Users/edmundd/Desktop/logo.gif"

    all_tags = get_all_tags_in_presentation(presentation)
    print "All tags: " + str(all_tags)
    total_slides = len(presentation.slides)

    if ARGS.trim_logo:
        trim(Image.open(logo_image)).save(TEMP_LOGO_FILE)
    else:
        Image.open(logo_image).save(TEMP_LOGO_FILE)

    for index, slide in enumerate(presentation.slides):
        print "{0}/{1}".format(index + 1, total_slides)

        update_logo(slide)

        tags_to_delete = ARGS.tags.split(",")

        if slide.has_notes_slide:
            text_frame = slide.notes_slide.notes_text_frame
            tags_in_comment = get_all_tags_in_comment(text_frame.text)
            matching_tags = copy.copy(tags_in_comment)
            if should_delete_slide(tags_in_comment, tags_to_delete):
                print "Deleting slide {0} with matching tags '{1}'".format(index, matching_tags)
                delete_slide(presentation, slide)
            text_frame.text = remove_json_from_comment(text_frame.text)

    os.remove(TEMP_LOGO_FILE)

    presentation.save(presentation_file.replace(".ppt", "-new.ppt"))
    print "Done!"


def should_delete_slide(tags_in_comment, tags_to_delete):
    if tags_in_comment != []:
        for tag in tags_to_delete:
            if tag in tags_in_comment:
                tags_in_comment.remove(tag)
        return tags_in_comment == []
    return False


def update_logo(slide):
    for shape in slide.placeholders:
        if shape.name == "Picture Placeholder 3":
            idx = shape.placeholder_format.idx
            slide.shapes.add_picture(TEMP_LOGO_FILE,
                                     slide.placeholders[idx].left,
                                     slide.placeholders[idx].top,
                                     None,
                                     slide.placeholders[idx].height)


if __name__ == '__main__':
    main()
